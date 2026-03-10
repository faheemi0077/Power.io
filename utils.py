from pandas import *
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from collections import Counter
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL
from pptx import Presentation
import os
import json
import urllib.request
import urllib.error

#raw chart generation
#generating the raw chart was succesful
def generate_chart(presentation, slide_number, chart_type, left=Inches(1), top=Inches(1.5), width=Inches(6), height=Inches(4)):
    """generates the chart itself"""
    #map of the charts (technically 1 functionality puts the chart on if its in the map)
    #interpret the chart type
    chart_type_map = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "clustered column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked column": XL_CHART_TYPE.COLUMN_STACKED,
    "100% stacked column": XL_CHART_TYPE.COLUMN_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line with markers": XL_CHART_TYPE.LINE_MARKERS,
    "stacked line": XL_CHART_TYPE.LINE_STACKED,
    "stacked line with markers": XL_CHART_TYPE.LINE_MARKERS_STACKED,
    "100% stacked line": XL_CHART_TYPE.LINE_STACKED_100,
    "100% stacked line with markers": XL_CHART_TYPE.LINE_MARKERS_STACKED_100,

    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,

    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "clustered bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked bar": XL_CHART_TYPE.BAR_STACKED,
    "100% stacked bar": XL_CHART_TYPE.BAR_STACKED_100,

    "area": XL_CHART_TYPE.AREA,
    "stacked area": XL_CHART_TYPE.AREA_STACKED,
    "100% stacked area": XL_CHART_TYPE.AREA_STACKED_100,

    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "xy scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter with smooth lines": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    "scatter with straight lines": XL_CHART_TYPE.XY_SCATTER_LINES,

    "radar": XL_CHART_TYPE.RADAR,
    "radar with markers": XL_CHART_TYPE.RADAR_MARKERS,
    "filled radar": XL_CHART_TYPE.RADAR_FILLED
}
#plop the raw chart
    slide = presentation.slides[slide_number - 1]
    chart_type_enum = chart_type_map.get(chart_type.lower())
    if chart_type_enum is None:
        raise ValueError(f"Unknown chart type: {chart_type}")
    scatter_types = [
        XL_CHART_TYPE.XY_SCATTER,
        XL_CHART_TYPE.XY_SCATTER_LINES,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS
    ]
    if chart_type_enum in scatter_types:
        chart_data = XyChartData()
        series = chart_data.add_series('Series 1')
        series.add_data_point(0, 0)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = ['Category 1']
        chart_data.add_series('Series 1', [0])
    graphic_frame = slide.shapes.add_chart(
        chart_type_enum,
        left, top, width, height,
        chart_data
    )
    chart = graphic_frame.chart
    chart._slide_number = slide_number
    if chart.has_title:
        try:
            chart.title.text_frame.clear()
        except Exception:
            pass
    chart.has_title = False
    return chart

def parse_excel(excel_path):
    """takes excel path, parses it for interpret_data to use"""
    if not isinstance(excel_path, str) or not excel_path.strip():
        raise ValueError("excel_path must be a non-empty string")

    if not os.path.exists(excel_path):
        raise FileNotFoundError(f"Excel file not found: {excel_path}")

    try:
        # Load the first worksheet only, returning a raw DataFrame as-is (pandas defaults).
        return read_excel(excel_path, sheet_name=0)
    except Exception as e:
        raise RuntimeError(f"Failed to read Excel file at '{excel_path}': {e}") from e
        
def interpret_data(dataframe, chart_type):
    """takes the data, passes it to the model, and returns a perfectly formatted data object for add_data to use"""
    # Serialize the raw DataFrame into a JSON-compatible structure (headers + rows).
    def _strip_json_fences(text: str) -> str:
        if not isinstance(text, str):
            return ""

        s = text.strip()

        if s.startswith("```"):
            # remove leading ```
            s = s.lstrip("`").strip()
            # remove optional 'json'
            if s.lower().startswith("json"):
                s = s[4:].strip()
            # remove trailing ```
            s = s.rstrip("`").strip()

        return s
    if dataframe is None:
        raise ValueError("dataframe must not be None")

    try:
        columns = [str(c) for c in list(dataframe.columns)]
        # JSON does not support NaN/Inf and json.dumps cannot serialize numpy scalar types reliably.
        values = dataframe.where(notna(dataframe), None).to_numpy().tolist()
        rows = []
        for r in values:
            out_r = []
            for v in r:
                if hasattr(v, "item") and callable(getattr(v, "item")):
                    try:
                        v = v.item()
                    except Exception:
                        pass
                out_r.append(v)
            rows.append(out_r)
        payload_table = {"columns": columns, "rows": rows}
    except Exception as e:
        raise RuntimeError(f"Failed to serialize dataframe for model: {e}") from e

    api_key = os.getenv("MISTRAL_API_KEY")
    if not api_key:
        raise RuntimeError("MISTRAL_API_KEY is not set")

    system_prompt = "You are a careful data reformatter for charting."
    user_prompt = (
        "You are given a raw table extracted from an Excel sheet.\n"
        "Your task is to reformat it into a clean table suitable for charting.\n\n"
        "Input:\n"
        f"- chart_type: {str(chart_type)}\n"
        "- table: JSON object with keys {columns: [..], rows: [[..], ..]}\n\n"
        "Output requirements (VERY IMPORTANT):\n"
        "- Return ONLY valid JSON. No prose. No markdown.\n"
        "- Output must be a JSON object with keys {columns: [...], rows: [[...], ...]} and may optionally include a key \"value_format\".\n"
        "- If present, \"value_format\" MUST be one of: \"number\", \"currency\", or \"percentage\".\n"
        "- Infer the semantic meaning of numeric values and set \"value_format\" accordingly.\n"
        "- The output table must have:\n"
        "  - Column 0: categories (strings)\n"
        "  - Column 1..N: numeric series (floats)\n"
        "- Keep numeric values numeric. Do NOT format numbers as strings.\n"
        "- No empty columns.\n"
        "- No totals, headers, notes, footers, or commentary rows.\n"
        "- Preserve the original row ordering of the meaningful data.\n"
        "- If chart_type is pie/donut, return EXACTLY 2 columns: [categories, values].\n\n"
        "Here is the raw input table JSON:\n"
        f"{json.dumps(payload_table, ensure_ascii=False)}\n\n"
        "Return ONLY the JSON object."
        "If category labels represent monetary ranges (e.g. 30k-50k, <30k, >150k), normalize them to include dollar signs (e.g. $30k–$50k, <$30k, >$150k)."
        "Do not alter category labels unless they clearly represent money ranges."
    )

    req_payload = {
        "model": "mistral-small",
        "temperature": 0,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    }

    req = urllib.request.Request(
        "https://api.mistral.ai/v1/chat/completions",
        data=json.dumps(req_payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=45) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        raise RuntimeError(f"Mistral API call failed: {e}") from e

    try:
        resp_data = json.loads(raw)
        content = resp_data["choices"][0]["message"]["content"]
    except Exception as e:
        raise RuntimeError(f"Unexpected Mistral response format: {e}") from e

    try:
        cleaned = json.loads(_strip_json_fences(content))
    except Exception as e:
        raise ValueError(f"Model did not return valid JSON: {e}\nRaw content:\n{content}") from e

    if not isinstance(cleaned, dict) or "columns" not in cleaned or "rows" not in cleaned:
        raise ValueError("Model JSON must be an object with keys 'columns' and 'rows'")

    df_clean = DataFrame(cleaned["rows"], columns=cleaned["columns"])
    if isinstance(cleaned, dict) and "value_format" in cleaned:
        df_clean.attrs["value_format"] = cleaned.get("value_format")
    return df_clean

#plop the data
def add_data(chart, dataframe):
    """
    Adds data to a chart from a pandas DataFrame.
    Takes chart object and dataframe directly (no file path needed).
    """
    chart_type_id = chart.chart_type
    scatter_types = [
        XL_CHART_TYPE.XY_SCATTER,
        XL_CHART_TYPE.XY_SCATTER_LINES,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS
    ]


    #directly adds the data if its possible
    if chart_type_id in scatter_types:
        chart_data = XyChartData()
        if len(dataframe.columns) < 2:
            raise ValueError("Scatter plot requires at least 2 columns (X and Y)")
        x_values = dataframe.iloc[:, 0].tolist()
        for col_idx in range(1, len(dataframe.columns)):
            series_name = str(dataframe.columns[col_idx])
            y_values = dataframe.iloc[:, col_idx].tolist()
            series = chart_data.add_series(series_name)
            for x, y in zip(x_values, y_values):
                series.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()

        if len(dataframe.columns) < 2:
            raise ValueError("Chart requires at least 2 columns (categories and values)")

        # categories always come from the first column
        categories = dataframe.iloc[:, 0].astype(str).tolist()
        chart_data.categories = categories

        # PIE / DOUGHNUT: exactly one numeric series, ignore empty columns
        if chart_type_id in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
            value_series = None

            for col_idx in range(1, len(dataframe.columns)):
                col = dataframe.iloc[:, col_idx]

                # skip completely empty columns
                if col.isna().all():
                    continue

                try:
                    value_series = col.astype(float).tolist()
                    break
                except Exception:
                    continue

            if value_series is None:
                raise ValueError("Pie chart requires at least one numeric column")

            chart_data.add_series("Values", value_series)

        # ALL OTHER CATEGORY CHARTS (bar, column, line, area, etc.)
        else:
            for col_idx in range(1, len(dataframe.columns)):
                col = dataframe.iloc[:, col_idx]

                # skip empty columns
                if col.isna().all():
                    continue

                series_name = str(dataframe.columns[col_idx])
                numeric_values = []

                for val in col.tolist():
                    try:
                        numeric_values.append(float(val))
                    except (ValueError, TypeError):
                        numeric_values.append(0)

                chart_data.add_series(series_name, numeric_values)
    chart.replace_data(chart_data)
    return chart



#styling 
#count theme + interpreting + applying the style seems to be working for now
def count_theme(presentation):
    """
    Analyze a PowerPoint Presentation and return a deck-level, JSON-serializable theme summary.

    Notes:
    - Builds per-slide theme profiles internally (not returned).
    - Only measured + aggregated data. No interpretation/styling decisions.
    - Skips unresolved theme defaults that cannot be converted to concrete RGB.
    - Does not mutate the presentation.
    """
    #functionality: counting the theme for the model to use
    # --- Deck-level aggregators ---
    slides_analyzed = 0

    # Colors
    bg_color_counts = Counter()
    color_counts = Counter()  # includes fills/outlines/text/chart colors; hex strings only

    # Typography
    font_counts = Counter()
    font_size_raw_counts = Counter()  # numeric point sizes as strings for JSON keys
    font_size_bucket_counts = Counter()
    bold_count = 0
    italic_count = 0
    underline_count = 0
    uppercase_count = 0
    total_text_runs = 0

    # Layout & density
    shapes_per_slide = []
    shape_type_counts = Counter()  # text_boxes, rectangles, lines, images, charts
    slides_with_images = 0
    slides_with_charts = 0
    slides_with_titles = 0

    # Charts
    chart_type_counts = Counter()
    chart_gridlines_counts = Counter()  # enabled/disabled
    chart_legend_position_counts = Counter()  # position name or 'none'
    chart_data_labels_counts = Counter()  # enabled/disabled

    # Internal only: slide-level profiles (NOT returned)
    slide_profiles = []

    for slide_idx, slide in enumerate(getattr(presentation, "slides", []), start=1):
        slides_analyzed += 1

        # --- Per-slide profile (internal) ---
        slide_profile = {
            "slide_index": slide_idx,
            "colors": {"background": Counter(), "fills": Counter(), "outlines": Counter(), "text": Counter(), "charts": Counter()},
            "typography": {
                "fonts": Counter(),
                "font_sizes": Counter(),
                "bold": 0,
                "italic": 0,
                "underline": 0,
                "uppercase": 0,
                "text_runs": 0,
            },
            "layout": {
                "total_shapes": 0,
                "shape_types": Counter(),
                "contains_image": False,
                "contains_chart": False,
            },
            "charts": {
                "chart_types": Counter(),
                "gridlines": Counter(),
                "legend_positions": Counter(),
                "data_labels": Counter(),
            },
            "metadata": {
                "has_chart": False,
                "has_image": False,
                "has_title": False,
                "slide_type_guess": "content_slide",
            },
        }

        slide_shapes = list(getattr(slide, "shapes", []))
        slide_shape_count = len(slide_shapes)
        shapes_per_slide.append(slide_shape_count)
        slide_profile["layout"]["total_shapes"] = slide_shape_count

        # Background color (solid only)
        bg_hex = None
        try:
            bg_fill = slide.background.fill
            if bg_fill and bg_fill.type == MSO_FILL.SOLID:
                bg_hex = _rgb_to_hex(bg_fill.fore_color.rgb)
        except Exception:
            bg_hex = None
        if bg_hex:
            slide_profile["colors"]["background"][bg_hex] += 1
            bg_color_counts[bg_hex] += 1
            color_counts[bg_hex] += 1  # include in overall palette map

        # Slide metadata flags
        title_shape = _safe_slide_title(slide)
        has_title = bool(title_shape is not None)
        if has_title:
            slides_with_titles += 1
        slide_profile["metadata"]["has_title"] = has_title

        has_image = False
        has_chart = False

        # Per-shape scanning
        for shape in slide_shapes:
            # Identify charts early (python-pptx uses GRAPHIC_FRAME for charts)
            if getattr(shape, "has_chart", False):
                has_chart = True
                slide_profile["layout"]["contains_chart"] = True
                slide_profile["metadata"]["has_chart"] = True
                slide_profile["layout"]["shape_types"]["charts"] += 1
                shape_type_counts["charts"] += 1

                # Chart style clues + chart colors (concrete only)
                try:
                    chart = shape.chart
                except Exception:
                    chart = None
                if chart is not None:
                    # Chart type
                    try:
                        ct = chart.chart_type
                        ct_name = getattr(ct, "name", str(ct))
                        slide_profile["charts"]["chart_types"][ct_name] += 1
                        chart_type_counts[ct_name] += 1
                    except Exception:
                        pass

                    # Gridlines enabled/disabled (best-effort; if unknown, omit)
                    gridlines_status = None
                    for axis_attr in ("category_axis", "value_axis"):
                        try:
                            axis = getattr(chart, axis_attr)
                        except Exception:
                            axis = None
                        if axis is None:
                            continue
                        try:
                            if hasattr(axis, "has_major_gridlines"):
                                gridlines_status = "enabled" if axis.has_major_gridlines else "disabled"
                                break
                        except Exception:
                            pass
                        try:
                            gridlines_status = "enabled" if axis.major_gridlines is not None else "disabled"
                            break
                        except Exception:
                            pass
                    if gridlines_status:
                        slide_profile["charts"]["gridlines"][gridlines_status] += 1
                        chart_gridlines_counts[gridlines_status] += 1

                    # Legend position
                    try:
                        if chart.has_legend:
                            pos = getattr(chart.legend.position, "name", str(chart.legend.position))
                        else:
                            pos = "none"
                        slide_profile["charts"]["legend_positions"][pos] += 1
                        chart_legend_position_counts[pos] += 1
                    except Exception:
                        pass

                    # Data labels enabled/disabled (best-effort; if unknown, omit)
                    try:
                        enabled = False
                        for s in chart.series:
                            try:
                                if hasattr(s, "has_data_labels") and s.has_data_labels:
                                    enabled = True
                                    break
                            except Exception:
                                continue
                        status = "enabled" if enabled else "disabled"
                        slide_profile["charts"]["data_labels"][status] += 1
                        chart_data_labels_counts[status] += 1
                    except Exception:
                        pass

                    # Chart colors (series fills; concrete RGB only)
                    try:
                        for s in chart.series:
                            try:
                                fill = s.format.fill
                                if fill is not None and fill.type == MSO_FILL.SOLID:
                                    hex_c = _rgb_to_hex(fill.fore_color.rgb)
                                    if hex_c:
                                        slide_profile["colors"]["charts"][hex_c] += 1
                                        color_counts[hex_c] += 1
                            except Exception:
                                continue
                    except Exception:
                        pass

                continue  # do not double-count chart as another shape category

            # Images
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                slide_profile["layout"]["contains_image"] = True
                slide_profile["metadata"]["has_image"] = True
                slide_profile["layout"]["shape_types"]["images"] += 1
                shape_type_counts["images"] += 1

            # Lines
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
                slide_profile["layout"]["shape_types"]["lines"] += 1
                shape_type_counts["lines"] += 1

            # Rectangles (best-effort)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if shape.auto_shape_type in {
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.SNIP_SINGLE_CORNER_RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.SNIP_SAME_SIDE_CORNER_RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.SNIP_DIAGONAL_CORNER_RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_SAME_SIDE_CORNER_RECTANGLE,
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_DIAGONAL_CORNER_RECTANGLE,
                    }:
                        slide_profile["layout"]["shape_types"]["rectangles"] += 1
                        shape_type_counts["rectangles"] += 1
                except Exception:
                    pass

            # Text boxes (treat any shape with a text frame as text-bearing element)
            if getattr(shape, "has_text_frame", False):
                slide_profile["layout"]["shape_types"]["text_boxes"] += 1
                shape_type_counts["text_boxes"] += 1

            # Shape fill colors (solid only; concrete RGB only)
            try:
                fill = shape.fill
                if fill is not None and fill.type == MSO_FILL.SOLID:
                    hex_c = _rgb_to_hex(fill.fore_color.rgb)
                    if hex_c:
                        slide_profile["colors"]["fills"][hex_c] += 1
                        color_counts[hex_c] += 1
            except Exception:
                pass

            # Shape outline/border colors (concrete RGB only)
            try:
                line = shape.line
                if line is not None and line.color is not None:
                    hex_c = _rgb_to_hex(line.color.rgb)
                    if hex_c:
                        slide_profile["colors"]["outlines"][hex_c] += 1
                        color_counts[hex_c] += 1
            except Exception:
                pass

            # Text run parsing
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            total_text_runs += 1
                            slide_profile["typography"]["text_runs"] += 1

                            # Font family
                            try:
                                if run.font.name:
                                    font_counts[run.font.name] += 1
                                    slide_profile["typography"]["fonts"][run.font.name] += 1
                            except Exception:
                                pass

                            # Font size (points)
                            try:
                                if run.font.size is not None:
                                    # Pt is an Emu wrapper; float() yields points for Pt objects
                                    pt_val = float(run.font.size.pt) if hasattr(run.font.size, "pt") else float(run.font.size)
                                    pt_val = float(pt_val)
                                    key = str(int(round(pt_val))) if pt_val.is_integer() else str(round(pt_val, 2))
                                    font_size_raw_counts[key] += 1
                                    slide_profile["typography"]["font_sizes"][key] += 1
                                    bucket = _bucket_font_size(pt_val)
                                    font_size_bucket_counts[bucket] += 1
                            except Exception:
                                pass

                            # Style flags
                            try:
                                if run.font.bold:
                                    bold_count += 1
                                    slide_profile["typography"]["bold"] += 1
                            except Exception:
                                pass
                            try:
                                if run.font.italic:
                                    italic_count += 1
                                    slide_profile["typography"]["italic"] += 1
                            except Exception:
                                pass
                            try:
                                if run.font.underline:
                                    underline_count += 1
                                    slide_profile["typography"]["underline"] += 1
                            except Exception:
                                pass

                            # Uppercase (optional)
                            try:
                                t = run.text or ""
                                if any(ch.isalpha() for ch in t) and t.upper() == t and t.lower() != t:
                                    uppercase_count += 1
                                    slide_profile["typography"]["uppercase"] += 1
                            except Exception:
                                pass

                            # Text run colors (concrete RGB only)
                            try:
                                hex_c = _rgb_to_hex(run.font.color.rgb)
                                if hex_c:
                                    slide_profile["colors"]["text"][hex_c] += 1
                                    color_counts[hex_c] += 1
                            except Exception:
                                pass
                except Exception:
                    pass

        if has_image:
            slides_with_images += 1
        if has_chart:
            slides_with_charts += 1

        slide_profile["metadata"]["has_chart"] = has_chart
        slide_profile["metadata"]["has_image"] = has_image
        slide_profile["metadata"]["has_title"] = has_title
        slide_profile["layout"]["contains_image"] = has_image
        slide_profile["layout"]["contains_chart"] = has_chart
        slide_profile["metadata"]["slide_type_guess"] = _guess_slide_type(
            slide, slide_shape_count, has_chart=has_chart, has_image=has_image, has_title=has_title
        )

        slide_profiles.append(slide_profile)

    # --- Aggregation into deck-level summary (returned) ---
    def _counter_to_dict(counter: Counter) -> dict:
        return {k: int(v) for k, v in counter.items() if v}

    def _most_common_key(counter: Counter) -> str | None:
        if not counter:
            return None
        return counter.most_common(1)[0][0]

    # Dominant background & ratios
    dominant_bg = _most_common_key(bg_color_counts)
    bg_consistency_ratio = 0.0
    if slides_analyzed and bg_color_counts:
        bg_consistency_ratio = float(bg_color_counts.most_common(1)[0][1]) / float(slides_analyzed)

    # Fonts & ratios
    primary_font = _most_common_key(font_counts)
    secondary_font = None
    if len(font_counts) >= 2:
        secondary_font = font_counts.most_common(2)[1][0]
    font_consistency_ratio = 0.0
    if total_text_runs and font_counts:
        font_consistency_ratio = float(font_counts.most_common(1)[0][1]) / float(sum(font_counts.values()))

    # Layout stats
    avg_shapes = float(sum(shapes_per_slide) / slides_analyzed) if slides_analyzed else 0.0
    median_shapes = 0.0
    if shapes_per_slide:
        s = sorted(shapes_per_slide)
        mid = len(s) // 2
        if len(s) % 2 == 1:
            median_shapes = float(s[mid])
        else:
            median_shapes = float((s[mid - 1] + s[mid]) / 2)

    # Density classification (based on median shapes per slide)
    if median_shapes <= 5:
        density = "low"
    elif median_shapes <= 12:
        density = "medium"
    else:
        density = "high"

    # Accent candidates: top non-background colors
    accent_candidates = []
    if color_counts:
        for c, _ in color_counts.most_common():
            if dominant_bg and c == dominant_bg:
                continue
            accent_candidates.append(c)
            if len(accent_candidates) >= 6:
                break

    # Color palette size
    color_palette_size = len([k for k, v in color_counts.items() if v])

    # Build response, omitting empty subcategories
    result = {
        "slides_analyzed": int(slides_analyzed),
        "colors": {},
        "fonts": {},
        "font_sizes": {},
        "layout": {},
        "charts": {},
        "consistency": {},
    }

    # Colors
    colors_section = {}
    if dominant_bg:
        colors_section["dominant_background_color"] = dominant_bg
    if color_counts:
        colors_section["color_frequency"] = _counter_to_dict(color_counts)
        colors_section["accent_color_candidates"] = accent_candidates
    if bg_color_counts:
        colors_section["background_color_frequency"] = _counter_to_dict(bg_color_counts)
    if colors_section:
        result["colors"] = colors_section

    # Fonts
    fonts_section = {}
    if primary_font:
        fonts_section["primary_font"] = primary_font
    if secondary_font and secondary_font != primary_font:
        fonts_section["secondary_font"] = secondary_font
    if font_counts:
        fonts_section["font_frequency"] = _counter_to_dict(font_counts)
    if font_counts:
        fonts_section["font_consistency_ratio"] = float(font_consistency_ratio)
    if fonts_section:
        result["fonts"] = fonts_section

    # Font sizes
    font_sizes_section = {}
    if font_size_bucket_counts:
        # Ensure stable keys
        font_sizes_section["bucketed_counts"] = {
            "large": int(font_size_bucket_counts.get("large", 0)),
            "medium": int(font_size_bucket_counts.get("medium", 0)),
            "small": int(font_size_bucket_counts.get("small", 0)),
        }
    if font_size_raw_counts:
        font_sizes_section["raw_point_size_frequency"] = _counter_to_dict(font_size_raw_counts)
    if font_sizes_section:
        result["font_sizes"] = font_sizes_section

    # Layout
    layout_section = {
        "average_shapes_per_slide": float(avg_shapes),
        "median_shapes_per_slide": float(median_shapes),
        "overall_density": density,
    }
    if shape_type_counts:
        layout_section["shape_type_counts"] = _counter_to_dict(shape_type_counts)
    layout_section["slides_with_titles"] = int(slides_with_titles)
    layout_section["slides_with_images"] = int(slides_with_images)
    layout_section["slides_with_charts"] = int(slides_with_charts)
    result["layout"] = layout_section

    # Charts
    charts_section = {"charts_exist": bool(slides_with_charts > 0)}
    if chart_type_counts:
        charts_section["chart_type_frequency"] = _counter_to_dict(chart_type_counts)
    if chart_gridlines_counts:
        charts_section["gridlines_frequency"] = _counter_to_dict(chart_gridlines_counts)
    if chart_legend_position_counts:
        charts_section["legend_position_frequency"] = _counter_to_dict(chart_legend_position_counts)
    if chart_data_labels_counts:
        charts_section["data_labels_frequency"] = _counter_to_dict(chart_data_labels_counts)
    result["charts"] = charts_section

    # Consistency
    uppercase_ratio = float(uppercase_count / total_text_runs) if total_text_runs else 0.0
    consistency_section = {
        "background_consistency_ratio": float(bg_consistency_ratio),
        "font_consistency_ratio": float(font_consistency_ratio),
        "color_palette_size": int(color_palette_size),
        "uppercase_ratio": float(uppercase_ratio),
    }
    result["consistency"] = consistency_section
    # Optional typography totals (only if present)
    if any([bold_count, italic_count, underline_count, uppercase_count]):
        result.setdefault("fonts", {})
        if bold_count:
            result["fonts"]["bold_text_count"] = int(bold_count)
        if italic_count:
            result["fonts"]["italic_text_count"] = int(italic_count)
        if underline_count:
            result["fonts"]["underline_text_count"] = int(underline_count)
        if uppercase_count:
            result["fonts"]["uppercase_text_count"] = int(uppercase_count)

    return result

def get_interpretation(theme_summary: dict) -> dict:
    """
    Takes the output of count_theme(theme_summary) and calls the Mistral AI API to infer chart styling intent.

    Returns:
      - A strict Python dict matching the required schema.
      - If the API call fails or the AI response is invalid/malformed/missing keys, returns safe defaults.
    """
    #interprets the theme through the model for the chart to have a given style
    fallback = {
        "use_theme_fonts": True,
        "use_theme_colors": True,
        "chart_title": {"enabled": False, "case": "sentence"},
        "axis_labels": {"case": "sentence"},
        "gridlines": False,
        "data_labels": False,
        "legend": {"enabled": True, "position": "bottom"},
        "visual_density": "minimal",
    }

    def _is_bool(x) -> bool:
        return isinstance(x, bool)

    def _validate_interpretation_dict(d: dict) -> bool:
        if not isinstance(d, dict):
            return False

        required = {
            "use_theme_fonts",
            "use_theme_colors",
            "chart_title",
            "axis_labels",
            "gridlines",
            "data_labels",
            "legend",
            "visual_density",
        }
        if any(k not in d for k in required):
            return False

        if not _is_bool(d["use_theme_fonts"]) or not _is_bool(d["use_theme_colors"]):
            return False
        if not _is_bool(d["gridlines"]) or not _is_bool(d["data_labels"]):
            return False

        if not isinstance(d["chart_title"], dict):
            return False
        if "enabled" not in d["chart_title"] or "case" not in d["chart_title"]:
            return False
        if not _is_bool(d["chart_title"]["enabled"]):
            return False
        if d["chart_title"]["case"] not in {"sentence", "upper"}:
            return False

        if not isinstance(d["axis_labels"], dict):
            return False
        if "case" not in d["axis_labels"]:
            return False
        if d["axis_labels"]["case"] not in {"sentence", "upper"}:
            return False

        if not isinstance(d["legend"], dict):
            return False
        if "enabled" not in d["legend"] or "position" not in d["legend"]:
            return False
        if not _is_bool(d["legend"]["enabled"]):
            return False
        if d["legend"]["position"] not in {"bottom", "right", "left", "top", "none"}:
            return False

        if d["visual_density"] not in {"minimal", "balanced", "dense"}:
            return False

        return True

    def _extract_json_object(text: str) -> dict | None:
        """
        AI is instructed to return only JSON, but we still parse defensively.
        """
        if not isinstance(text, str):
            return None
        s = text.strip()
        # Prefer strict whole-string JSON
        try:
            obj = json.loads(s)
            return obj if isinstance(obj, dict) else None
        except Exception:
            pass
        # Best-effort: extract first {...} block
        start = s.find("{")
        end = s.rfind("}")
        if start == -1 or end == -1 or end <= start:
            return None
        try:
            obj = json.loads(s[start : end + 1])
            return obj if isinstance(obj, dict) else None
        except Exception:
            return None

    api_key = os.getenv("MISTRAL_API_KEY")
    if not api_key:
        return fallback

    # Serialize theme_summary and inject into prompt.
    try:
        theme_json = json.dumps(theme_summary, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    except Exception:
        theme_json = json.dumps({"theme_summary": str(theme_summary)}, ensure_ascii=False, sort_keys=True, separators=(",", ":"))

    system_prompt = "You are an expert PowerPoint chart styling assistant."

    user_prompt = (
        "You are given a JSON object containing ONLY measured, aggregated theme attributes\n"
        "from an existing PowerPoint presentation.\n\n"
        "Your task is to infer chart styling intent so that newly created charts visually\n"
        "match the presentation.\n\n"
        "Rules:\n"
        "- Do NOT guess brand colors or fonts.\n"
        "- If font or color information is missing, inherit from the PowerPoint theme.\n"
        "- Prefer minimal, professional styling.\n"
        "- Match existing chart behavior when signals are present.\n"
        "- Never invent information.\n"
        "- Return ONLY valid JSON. No prose. No markdown.\n\n"
        "Return JSON with the following schema:\n\n"
        "{\n"
        '  "use_theme_fonts": boolean,\n'
        '  "use_theme_colors": boolean,\n\n'
        '  "chart_title": {\n'
        '    "enabled": boolean,\n'
        '    "case": "sentence" | "upper"\n'
        "  },\n\n"
        '  "axis_labels": {\n'
        '    "case": "sentence" | "upper"\n'
        "  },\n\n"
        '  "gridlines": boolean,\n'
        '  "data_labels": boolean,\n\n'
        '  "legend": {\n'
        '    "enabled": boolean,\n'
        '    "position": "bottom" | "right" | "left" | "top" | "none"\n'
        "  },\n\n"
        '  "visual_density": "minimal" | "balanced" | "dense"\n'
        "}\n\n"
        "Here is the measured theme summary:\n"
        f"{theme_json}\n\n"
        "Return ONLY the JSON object."
    )

    payload = {
        "model": "mistral-small",
        "temperature": 0,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    }

    req = urllib.request.Request(
        "https://api.mistral.ai/v1/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=25) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
    except Exception:
        return fallback

    try:
        data = json.loads(raw)
        content = data["choices"][0]["message"]["content"]
    except Exception:
        return fallback

    parsed = _extract_json_object(content)
    if parsed is None or not _validate_interpretation_dict(parsed):
        return fallback

    return parsed

def style_chart(chart, interpretation: dict) -> None:
    #purely styles the chart based on what the model interpreted
    #should the chart have a legend or not based on the chart type
    chart_type_id = chart.chart_type

    no_axis_chart_types = {
        XL_CHART_TYPE.PIE,
        XL_CHART_TYPE.DOUGHNUT,
        XL_CHART_TYPE.RADAR,
        XL_CHART_TYPE.RADAR_MARKERS,
        XL_CHART_TYPE.RADAR_FILLED,
    }

    force_legend = chart_type_id in no_axis_chart_types
    # 1) Chart title
    try:
        if isinstance(interpretation, dict) and isinstance(interpretation.get("chart_title"), dict):
            enabled = interpretation["chart_title"].get("enabled")
            if enabled is False:
                chart.has_title = False
            elif enabled is True:
                chart.has_title = True
    except Exception:
        pass

    # 2) Axis label casing (only if axis titles already exist)
    try:
        if isinstance(interpretation, dict) and isinstance(interpretation.get("axis_labels"), dict):
            if interpretation["axis_labels"].get("case") == "upper":
                for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
                    if axis is not None and axis.has_title:
                        tf = axis.axis_title.text_frame
                        if tf.text is not None:
                            tf.text = tf.text.upper()
    except Exception:
        pass

    # 3) Gridlines (category + value axes if present)
    try:
        if isinstance(interpretation, dict) and isinstance(interpretation.get("gridlines"), bool):
            val = interpretation["gridlines"]
            for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
                if axis is not None:
                    axis.has_major_gridlines = val
    except Exception:
        pass


    # 4b) Data label number formatting (only if value_format is present)
    try:
        if isinstance(interpretation, dict) and "value_format" in interpretation:
            vf = interpretation.get("value_format")
            fmt = None
            if vf == "currency":
                fmt = "$#,##0"
            elif vf == "percentage":
                fmt = "0%"
            elif vf == "number":
                fmt = "#,##0"
            if fmt:
                for series in chart.series:
                    dl = series.data_labels
                    dl.number_format = fmt
                    dl.number_format_is_linked = False
    except Exception:
        pass

    # 5) Legend
    single_series = len(chart.series) == 1
    if force_legend:
        chart.has_legend = True
        chart.legend.font.size = Pt(7)

    else:
        # Axis charts: hide legend for single-series charts regardless of AI interpretation
        if single_series:
            chart.has_legend = False
        else:
            if isinstance(interpretation, dict) and isinstance(interpretation.get("legend"), dict):
                enabled = interpretation["legend"].get("enabled", True)
                position = interpretation["legend"].get("position", "bottom")

                chart.has_legend = bool(enabled)

                if chart.has_legend:
                    chart.legend.font.size = Pt(7)
                    if position == "bottom":
                        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    elif position == "top":
                        chart.legend.position = XL_LEGEND_POSITION.TOP
                    elif position == "left":
                        chart.legend.position = XL_LEGEND_POSITION.LEFT
                    elif position == "right":
                        chart.legend.position = XL_LEGEND_POSITION.RIGHT

def get_special_instructions(special_instructions=None):
    """
    Retrieve raw user-provided special instructions for chart customization.

    This function is intentionally minimal:
    - No AI calls
    - No parsing / interpretation / validation
    - No styling logic, chart logic, or defaults
    - No side effects
    """
    if special_instructions is None or special_instructions == "":
        return None
    return special_instructions

def pass_special_instructions(special_instructions: str | None) -> dict | None:
    """
    Interpret raw user special instructions into a strictly validated override schema.

    - Calls the Mistral API (deterministic: temperature=0)
    - Returns ONLY the validated override dictionary, or None on any failure/invalidity
    - Does not apply any chart logic/styling and does not reference python-pptx objects
    """
    if special_instructions is None or not str(special_instructions).strip():
        return None

    api_key = os.getenv("MISTRAL_API_KEY")
    if not api_key:
        return None

    allowed_legend_positions = {"bottom", "top", "left", "right"}
    required_top_keys = {
        "chart_title",
        "axis_labels",
        "swap_axes",
        "gridlines",
        "data_labels",
        "legend",
    }

    def _is_bool_or_none(v) -> bool:
        return v is None or isinstance(v, bool)

    def _strip_json_fences(text: str) -> str:
        if not isinstance(text, str):
            return ""
        s = text.strip()
        if s.startswith("```"):
            s = s.lstrip("`").strip()
            if s.lower().startswith("json"):
                s = s[4:].strip()
            s = s.rstrip("`").strip()
        return s

    def _validate_schema(obj: dict) -> bool:
        if not isinstance(obj, dict):
            return False
        if set(obj.keys()) != required_top_keys:
            return False

        # chart_title / axis_labels
        for k in ("chart_title", "axis_labels"):
            v = obj.get(k)
            if not isinstance(v, dict) or set(v.keys()) != {"enabled"}:
                return False
            if not _is_bool_or_none(v.get("enabled")):
                return False

        # top-level toggles
        for k in ("swap_axes", "gridlines", "data_labels"):
            if not _is_bool_or_none(obj.get(k)):
                return False

        # legend
        legend = obj.get("legend")
        if not isinstance(legend, dict) or set(legend.keys()) != {"enabled", "position"}:
            return False
        if not _is_bool_or_none(legend.get("enabled")):
            return False
        pos = legend.get("position")
        if not (pos is None or (isinstance(pos, str) and pos in allowed_legend_positions)):
            return False

        return True

    system_prompt = (
        "You strictly extract explicit user intent into a predefined JSON schema. "
        "Never invent features. Never guess. Never infer. "
        "Return ONLY valid JSON with the exact required keys. No prose. No markdown."
    )

    user_prompt = (
        "You will receive a user's natural-language instructions for chart customization.\n\n"
        "Your job is NOT to explain anything and NOT to guess intent.\n"
        "Your ONLY job is to output a JSON object that STRICTLY matches the schema below.\n\n"
        "CRITICAL RULES (follow exactly):\n"
        "- You MUST output valid JSON.\n"
        "- You MUST include ALL keys shown in the schema.\n"
        "- You MUST NOT add extra keys.\n"
        "- You MUST NOT remove any keys.\n"
        "- You MUST NOT change nesting.\n"
        "- You MUST NOT include prose, markdown, comments, or explanations.\n"
        "- If the user explicitly mentions a feature, set the corresponding value to true or false.\n"
        "- If the user does NOT explicitly mention a feature, set it to null.\n"
        "- Do NOT infer intent. Do NOT assume defaults.\n"
        "- Unsupported requests MUST be ignored and represented as null.\n\n"
        "REQUIRED JSON SCHEMA (output EXACTLY this structure):\n"
        "{\n"
        '  "chart_title": { "enabled": true | false | null },\n'
        '  "axis_labels": { "enabled": true | false | null },\n'
        '  "swap_axes": true | false | null,\n'
        '  "gridlines": true | false | null,\n'
        '  "data_labels": true | false | null,\n'
        '  "legend": {\n'
        '    "enabled": true | false | null,\n'
        '    "position": "bottom" | "top" | "left" | "right" | null\n'
        "  }\n"
        "}\n\n"
        "USER INSTRUCTIONS (raw text):\n"
        f"{special_instructions.strip()}\n\n"
        "REMEMBER: Output ONLY the JSON object. Nothing else."
    )

    payload = {
        "model": "mistral-small",
        "temperature": 0,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    }

    req = urllib.request.Request(
        "https://api.mistral.ai/v1/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=25) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
    except Exception:
        return None

    try:
        data = json.loads(raw)
        content = data["choices"][0]["message"]["content"]
    except Exception:
        return None

    try:
        parsed = json.loads(_strip_json_fences(content))
    except Exception:
        return None

    # Normalize nullable nested fields returned as null by the model
    if parsed.get("chart_title") is None:
        parsed["chart_title"] = {"enabled": None}

    if parsed.get("axis_labels") is None:
        parsed["axis_labels"] = {"enabled": None}

    if not _validate_schema(parsed):
        return None

    return parsed

# def apply_special_instructions(chart, dataframe, interpretation, special_overrides):
#     """
#     Apply validated user override instructions on top of an existing interpretation.

#     - No AI calls
#     - Only toggles on/off features and performs axis swapping
#     - Defensive against unsupported python-pptx attributes / chart types
#     - Returns (chart, dataframe) (possibly updated)
#     """
#     if not special_overrides:
#         return chart, dataframe

#     overrides = special_overrides  # do not mutate interpretation

#     def _get(d, *path):
#         cur = d
#         for k in path:
#             if not isinstance(cur, dict) or k not in cur:
#                 return None
#             cur = cur[k]
#         return cur

#     # 1) Chart title
#     ct_enabled = _get(overrides, "chart_title", "enabled")
#     if ct_enabled is True:
#         try:
#             chart.has_title = True
#         except Exception:
#             pass
#     elif ct_enabled is False:
#         try:
#             chart.has_title = False
#         except Exception:
#             pass

#     # 2) Axis labels: only remove titles if explicitly disabled
#     al_enabled = _get(overrides, "axis_labels", "enabled")
#     if al_enabled is False:
#         for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
#             try:
#                 if axis is not None and getattr(axis, "has_title", False):
#                     axis.has_title = False
#             except Exception:
#                 pass

#     # 3) Gridlines
#     gridlines = _get(overrides, "gridlines")
#     if isinstance(gridlines, bool):
#         for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
#             try:
#                 if axis is not None:
#                     axis.has_major_gridlines = gridlines
#             except Exception:
#                 pass

#     # 4) Data labels
#     data_labels = _get(overrides, "data_labels")
#     if isinstance(data_labels, bool):
#         try:
#             for series in getattr(chart, "series", []):
#                 try:
#                     series.has_data_labels = data_labels
#                     if data_labels is True:
#                         try:
#                             series.data_labels.show_value = True
#                         except Exception:
#                             pass
#                 except Exception:
#                     pass
#         except Exception:
#             pass

#     # 5) Legend
#     legend_enabled = _get(overrides, "legend", "enabled")
#     if legend_enabled is False:
#         try:
#             chart.has_legend = False
#         except Exception:
#             pass
#     elif legend_enabled is True:
#         try:
#             chart.has_legend = True
#         except Exception:
#             pass

#     legend_position = _get(overrides, "legend", "position")
#     if isinstance(legend_position, str):
#         try:
#             if getattr(chart, "has_legend", False) is False:
#                 chart.has_legend = True
#             if legend_position == "bottom":
#                 chart.legend.position = XL_LEGEND_POSITION.BOTTOM
#             elif legend_position == "top":
#                 chart.legend.position = XL_LEGEND_POSITION.TOP
#             elif legend_position == "left":
#                 chart.legend.position = XL_LEGEND_POSITION.LEFT
#             elif legend_position == "right":
#                 chart.legend.position = XL_LEGEND_POSITION.RIGHT
#         except Exception:
#             pass

#     return chart, dataframe